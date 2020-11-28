import os
from pathlib import Path
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from datetime import datetime
#import matplotlib.pyplot as plt

mhz = 1e6
spurious_emission_limit = 10 # 10dB from noise floor
span = 20 * mhz
wifi_channels = [{
		'2G': {
				"Ch 1":	2412000000,
				"Ch 2": 2417000000,
				"Ch 3": 2422000000,
				"Ch 4": 2427000000,
				"Ch 5": 2432000000,
				"Ch 6": 2437000000,
				"Ch 7": 2442000000,
				"Ch 8": 2447000000,
				"Ch 9": 2452000000,
				"Ch 10": 2457000000,
				"Ch 11": 2462000000
			},
		'5G': {
				"Ch 36": 5180000000,
				"Ch 38": 5190000000,
				"Ch 40": 5200000000,
				"Ch 42": 5210000000,
				"Ch 44": 5220000000,
				"Ch 46": 5230000000,
				"Ch 48": 5240000000,
				"Ch 50": 5250000000,
				"Ch 52": 5260000000,
				"Ch 54": 5270000000,
				"Ch 56": 5280000000,
				"Ch 58": 5290000000,
				"Ch 60": 5300000000,
				"Ch 62": 5310000000,
				"Ch 64": 5320000000,
				"Ch 68": 5340000000,
				"Ch 96": 5480000000,
				"Ch 100": 5500000000,
				"Ch 102": 5510000000,
				"Ch 104": 5520000000,
				"Ch 106": 5530000000,
				"Ch 108": 5540000000,
				"Ch 110": 5550000000,
				"Ch 112": 5560000000,
				"Ch 114": 5570000000,
				"Ch 116": 5580000000,
				"Ch 118": 5590000000,
				"Ch 120": 5600000000,
				"Ch 122": 5610000000,
				"Ch 124": 5620000000,
				"Ch 126": 5630000000,
				"Ch 128": 5640000000,
				"Ch 130": 5650000000,
				"Ch 132": 5660000000,
				"Ch 134": 5670000000,
				"Ch 136": 5680000000,
				"Ch 138": 5690000000,
				"Ch 140": 5700000000,
				"Ch 142": 5710000000,
				"Ch 144": 5720000000,
				"Ch 149": 5745000000,
				"Ch 151": 5755000000,
				"Ch 153": 5765000000,
				"Ch 155": 5775000000,
				"Ch 157": 5785000000,
				"Ch 159": 5795000000,
				"Ch 161": 5805000000,
				"Ch 165": 5825000000
			}
		}]

class powerpoint:
	def ppt_init(name, slides):
		try:
			prs = Presentation()
			slide1 = prs.slides.add_slide(prs.slide_layouts[1])
			prs.save(ppt_file)
		except Exception as e:
			print("*"*80 + "\nError generating powerpoint.\n" + "*"*80)
			raise e


	def chart_to_ppt(chart_dict):
		try:
			print("Chart dict:\t\t{}".format(chart_dict))


			
		except Exception as e:
			print("*"*80 + "\nError creating chart in Powerpoint.\n" + "*"*80)
			raise e

def __chart_points(df, result_type, band): # todo: very low priority...need seperate data_series for markers
	points = list()
	try:
		header = df.columns.values
		df['{}Delta'.format(result_type)] = df['{}Uncalibrated[dB]'.format(result_type)] - df['Noise Floor']
		spur_df = df['Frequency[Hz]'].where(df['{}Delta'.format(result_type)] > spurious_emission_limit).dropna()

		for ch, f in wifi_channels[0][band].items():
			for freq in spur_df:
				if int(freq) in range(int(f - span / 2), int(f + span/2)):
					print("Found a spur at {}".format(ch))
					points.append({'fill': {'color': 'green'}})
				else:
					points.append({'fill': {'none': True}})
			points.append({'fill': {'none': True}})
			print(points)
		return points
	except Exception as e:
		print("*"*80 + "\nError processing chart markers.\n" + "*"*80)
		raise e

def __spurs_table(worksheet): # TODO
	spurs_dict = dict()
	updated_worksheet = worksheet
	try:
		print("bleh")


		return updated_worksheet
	except Exception as e:
		print("*"*80 + "\nError generating spurs table.\n" + "*"*80)
		raise e	

def __integrated_power(channel, data):
	# integrated/channel power
	try:
		print("bleh2")
	except Exception as e:
		print("*"*80 + "\nError calculating integrated power.\n" + "*"*80)
		raise e	 # TOFINISH

def __noise_floor(nf_folder):
	nf_dfs = dict()
	try:
		for file in os.listdir(nf_folder):
			if ".csv" in file:
				nf_dfs[file] = pd.read_csv(nf_folder + "\\" + file)
		return nf_dfs
	except Exception as e:
		print("*"*80 + "\nError handling noise floor data.\n" + "*"*80)
		raise e

def csv_to_dict(results_folder):
	print("Parsing {}".format(results_folder))
	dfs = dict() # DataFrame(s) of parsed results
	link_details = {
		'radio': [],
		'antenna': [],
		'band': []
	}
	try:
		noise_floor = __noise_floor(str(Path(results_folder).parent) + "\\Noise Floor")
		for file in os.listdir(results_folder):
			if ".csv" in file:
				split_file = file.split("__")
				if split_file[1] not in link_details['radio']:
					link_details['radio'].append(split_file[1])
				if split_file[2] not in link_details['antenna']:
					link_details['antenna'].append(split_file[2])
				if split_file[3] not in link_details['band']:
					link_details['band'].append(split_file[3])

				dfs["_".join(split_file[0:4])] = pd.read_csv(results_folder + "\\" + file)
				dfs["_".join(split_file[0:4])]['Noise Floor'] = noise_floor["_".join(split_file[1:4]) + ".csv"]['Noise Floor']
			else:
				raise RuntimeError("Wrong filetype found; omitting {}".format(file))
		return (dfs, results_folder.split("\\"), link_details)
	except Exception as e:
		print("*"*80 + "\nError parsing results folder.\n" + "*"*80)
		raise e

def excel_handler(dfs_tuple):
	try:
		with pd.ExcelWriter(os.getcwd() + "\\Results\\" + dfs_tuple[1][-1] + "_" + datetime.today().strftime('%Y%M%d%H%M%S') + ".xlsx", engine='xlsxwriter') as writer:
			print("Writing to spreadsheet...")
			avg_chart_index = 0
			max_chart_index = 0
			workbook = writer.book
			summary_max = workbook.add_worksheet("Summary_MaxHold")
			summary_avg = workbook.add_worksheet("Summary_Average")
			link_info = dfs_tuple[2]
			for radio in link_info['radio']:
				for antenna in link_info['antenna']:
					if "nw" in radio.lower() and ("pri" in antenna.lower() or "div" in antenna.lower()):
						continue
					elif "acc" in radio.lower() and ("ant1" in antenna.lower() or "ant2" in antenna.lower()):
						continue
					for band in link_info['band']:
						column_index = 1
						if "div" in antenna.lower() and "2g" in band.lower():
							continue
						workbook.add_worksheet("{} {} {}".format(radio, antenna, band))
						col_length = 0
						for i, result_type in enumerate(['MaxHold', 'Average']):
							noise_floor_added = False
							initial_cols = False
							chart = workbook.add_chart({
								'type': 	'scatter',
								'subtype': 	'straight'
							})
							chart.set_legend({'none': True})
							chart.set_y_axis({
								#'name':		"Log Mag (dBm)",
								'min': 		-125,
								'max': 		-90
							})
							for key, values in dfs_tuple[0].items():
								key = key.replace("_", "").replace(" ", "")
								if workbook.get_worksheet_by_name(key) is None:
									values.to_excel(writer, sheet_name=key, index=False)
								worksheet = workbook.get_worksheet_by_name(key)
								worksheet.hide()
							for key, values in dfs_tuple[0].items():
								test_case = key.split("_")[0]
								if key == (test_case + "_" + radio + "_" + antenna + "_" + band):
									parsing_key = key.replace("_", "").replace(" ", "")
									x_min = 4.8e9
									x_max = 6.2e9
									if "2G" in parsing_key:
										x_min = 2.1e9
										x_max = 2.8e9

									chart.set_x_axis({
										#'name': 			"Frequency (MHz)",
										'num_format': 		'#.#,,',
										'min': 				x_min,
										'max': 				x_max,
										'label_position': 	'low',
										'major_gridlines': 	{'visible': False}
									})

									chart.set_title({
										'name': "{} {} {}".format(radio, antenna, band)
									})
									column = "E"
									if "Average" in result_type:
										column = "D"

									chart.add_series({ # selected cells tied to analyzer (span/num_pts) state.
										'name':			test_case,
										'categories': 	'={}!$A2:$A$1002'.format(parsing_key), # check dataframe.index for stop+1 instead of 1002
										'values': 		'={}!${}2:${}$1002'.format(parsing_key, column, column),
										'line':			{'width': 0.75}
									})
									if noise_floor_added is False:
										chart.add_series({ # selected cells tied to analyzer (span/num_pts) state.
											'name': 		"Noise Floor",
											'categories': 	'={}!$A2:$A$1002'.format(parsing_key),
											'values': 		'={}!$G2:$G$1002'.format(parsing_key),
											'line':			{'width': 0.75}
										})
										noise_floor_added = True
									updated_worksheet = workbook.get_worksheet_by_name("{} {} {}".format(radio, antenna, band))
									parsed_test_case = key.replace("_" + radio + "_" + antenna + "_" + band, "")
									if initial_cols is False:
										if i == 1:
											column_index = 1
										for cols in ["Frequency[Hz]", "Noise Floor"]:
											row_index = len(values[cols]) * i
											if i == 0:
												col_length += len(values[cols])
												updated_worksheet.write_string(0 + row_index, column_index, cols)
												updated_worksheet.write_string(0 + row_index, 0, result_type)
											updated_worksheet.write_column(1 + row_index, column_index, values[cols])
											column_index += 1
										updated_worksheet.write_column(1 + row_index, 0, [str(result_type) for i in range(len(values[cols]))])
										initial_cols = True
									if i == 0:
										updated_worksheet.write_string(0 + row_index, column_index, parsed_test_case)
									updated_worksheet.write_column(1 + row_index, column_index, values["{}Uncalibrated[dB]".format(result_type)])
									column_index += 1
							if "Average" in result_type:
								summary_avg.insert_chart(avg_chart_index*10, 1, chart)
								avg_chart_index += 2
							else:
								summary_max.insert_chart(max_chart_index*10, 1, chart)
								#powerpoint.chart_to_ppt(chart.__dict__)
								max_chart_index += 2
						# TODO: add new columns to updated_worksheet with spurs?
						#updated_worksheet = __spurs_table(updated_worksheet) # TODO: pass updated_worksheet to func and back?
						updated_worksheet.add_table('A1:B{}'.format(col_length+1), {'columns': [{'header': 'MeasType'},
																								{'header': 'Frequency'}
																					]
																				})
			#print("Chart info: {}".format(chart.__dict__['series']))
	except Exception as e:
		print("*"*80 + "\nError writing spreadsheet.\n" + "*"*80)
		raise e

if __name__ == "__main__":
	#alias: v-brala
	time_start = datetime.now()
	dir_list = list()
	results_folder = os.getcwd() + "\\Results"
	for root, dirs, files in os.walk(results_folder):
		dir_list.append(dirs)
	for folders in dir_list[0]:
		if folders == "Noise Floor":
			continue
		parsed_data = csv_to_dict(results_folder + "\\" + folders)
		excel_handler(parsed_data)

	print("Complete!")
	print("Total processing time: {}s".format(datetime.now() - time_start))
